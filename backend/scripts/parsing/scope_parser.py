import fitz  # PyMuPDF
import docx
from pptx import Presentation


def parse_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""

    for page in doc:
        text += page.get_text()

    return text


def parse_docx(file_path):
    doc = docx.Document(file_path)

    text = "\n".join([p.text for p in doc.paragraphs])

    return text


def parse_pptx(file_path):
    prs = Presentation(file_path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def parse_scope_document(file_path):

    if file_path.endswith(".pdf"):
        return parse_pdf(file_path)

    elif file_path.endswith(".docx"):
        return parse_docx(file_path)

    elif file_path.endswith(".pptx"):
        return parse_pptx(file_path)

    else:
        print(f"Unsupported file type: {file_path}")
        return None